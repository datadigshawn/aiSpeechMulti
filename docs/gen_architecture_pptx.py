#!/usr/bin/env python3
"""
aiSpeechMulti 語音辨識平台架構說明簡報產生器
輸出: docs/aiSpeechMulti_Architecture.pptx

更新版本: v2.1 (2026-04-06)
- 同步至最新平台架構（10 頁 Streamlit、184 條詞彙、170 條修正規則）
- 新增 SenseVoice 即時本地端模式（2026-04-01）
- 新增離線模式規劃章節（Mac mini M4 + 雙引擎）
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ── 色彩定義（Midnight Executive 主題）──────────────────────────────────────
BG_DARK   = RGBColor(0x12, 0x16, 0x2B)   # 深夜藍（封面/結尾）
BG_LIGHT  = RGBColor(0xF5, 0xF7, 0xFA)   # 淺灰白（內容頁）
PRIMARY   = RGBColor(0x1E, 0x40, 0xAF)   # 寶藍
ACCENT    = RGBColor(0x06, 0xB6, 0xD4)   # 青色
ACCENT2   = RGBColor(0x10, 0xB9, 0x81)   # 綠色
ACCENT3   = RGBColor(0xF5, 0x9E, 0x0B)   # 琥珀色
DANGER    = RGBColor(0xEF, 0x44, 0x44)   # 紅色
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK  = RGBColor(0x1E, 0x29, 0x3B)
TEXT_MID   = RGBColor(0x64, 0x74, 0x8B)
CARD_BG    = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_dark_slide(prs, title_text, subtitle_text=""):
    """深色背景標題頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_DARK

    # 左側裝飾條
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # 標題
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xA0, 0xAE, 0xC0)
        p2.font.name = "Calibri"
        p2.space_before = Pt(16)

    return slide


def add_content_slide(prs, title_text):
    """淺色背景內容頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_LIGHT

    # 頂部色條
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    # 標題
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.font.name = "Calibri"

    return slide


def add_card(slide, x, y, w, h, fill_color=CARD_BG):
    """白色卡片（帶陰影效果用邊框模擬）"""
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(1)
    return card


def add_text_in_card(slide, x, y, w, h, lines, font_size=13):
    """在卡片區域內新增文字"""
    txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(w - 0.4), Inches(h - 0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)


def add_stat_box(slide, x, y, number, label, color):
    """大數字統計方塊"""
    add_card(slide, x, y, 2.5, 1.5)
    # 數字
    txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(2.1), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Calibri"
    # 標籤
    txBox2 = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.95), Inches(2.1), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_MID
    p2.font.name = "Calibri"


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ══════════════════════════════════════════════════════════════════════
    # Slide 1: 封面
    # ══════════════════════════════════════════════════════════════════════
    slide = add_dark_slide(prs,
        "aiSpeechMulti",
        "多引擎即時語音辨識平台\n架構與核心技術說明"
    )
    # 版本資訊
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(5.0), Inches(10), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "v2.1  |  2026-04-06  |  6 路即時串流 + 5 種辨識引擎 + 10 頁管理面板"
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT
    p.font.name = "Calibri"

    # ══════════════════════════════════════════════════════════════════════
    # Slide 2: 平台總覽
    # ══════════════════════════════════════════════════════════════════════
    slide = add_content_slide(prs, "平台總覽")

    # 統計方塊
    add_stat_box(slide, 0.8, 1.3, "6 路", "即時串流管道", PRIMARY)
    add_stat_box(slide, 3.6, 1.3, "5 種", "辨識引擎", ACCENT)
    add_stat_box(slide, 6.4, 1.3, "5 模式", "WebSocket 串流模式", ACCENT2)
    add_stat_box(slide, 9.2, 1.3, "10 頁", "Streamlit 管理面板", ACCENT3)

    # 核心能力卡片
    add_card(slide, 0.8, 3.2, 5.5, 3.8)
    add_text_in_card(slide, 0.8, 3.2, 5.5, 3.8, [
        ("核心能力", True, PRIMARY),
        ("", False, TEXT_MID),
        ("即時語音辨識 — 6 路 WebSocket 串流，~150ms 首字延遲", False, TEXT_DARK),
        ("批次音檔辨識 — 5 種引擎可選 + 6 步驟流程", False, TEXT_DARK),
        ("準確率評測 — CER/WER 字元級比對分析", False, TEXT_DARK),
        ("全文檢索 — SQLite FTS5 中文三元組搜尋", False, TEXT_DARK),
        ("詞彙優化 — PhraseSet + 術語修正 + 同音字修正", False, TEXT_DARK),
        ("事件管理 — 危害等級分類 + 音檔歸檔 + 關鍵字提取", False, TEXT_DARK),
        ("離線辨識 — SenseVoice / Whisper 100% 本地端", False, TEXT_DARK),
    ], font_size=13)

    add_card(slide, 6.8, 3.2, 5.5, 3.8)
    add_text_in_card(slide, 6.8, 3.2, 5.5, 3.8, [
        ("技術棧", True, PRIMARY),
        ("", False, TEXT_MID),
        ("後端: FastAPI + Streamlit + SQLite (WAL)", False, TEXT_DARK),
        ("前端: HTML5 + Web Audio API + WebSocket", False, TEXT_DARK),
        ("STT: Google chirp_3 / Scribe v2 / Whisper / SenseVoice / Gemini", False, TEXT_DARK),
        ("音訊: Silero VAD + DeepFilterNet 降噪 + ffmpeg 轉檔", False, TEXT_DARK),
        ("NLP: opencc 簡繁轉換 + jiwer CER/WER 評測", False, TEXT_DARK),
        ("ML: PyTorch + FunASR + faster-whisper (CTranslate2)", False, TEXT_DARK),
    ], font_size=13)

    # ══════════════════════════════════════════════════════════════════════
    # Slide 3: 系統架構圖
    # ══════════════════════════════════════════════════════════════════════
    slide = add_content_slide(prs, "系統架構（三層）")

    layers = [
        ("前端層", 1.3, ACCENT,  [
            "index.html — 單頻道串流控制 + 辨識結果顯示",
            "monitor.html — 6 路即時監控儀表板",
            "display.html — 全螢幕關鍵字投影（控制室）",
            "Web Audio API → PCM 16kHz Mono → WebSocket",
        ]),
        ("後端層", 3.2, PRIMARY, [
            "app_api.py — FastAPI（WebSocket + REST API + 6 路管道）",
            "app_dashboard.py — Streamlit 管理面板（10 頁面）",
            "SQLite WAL — aiSpeechMulti.db（即時 transcripts 表）",
            "AudioBuffer — 15s（batch）/ 3s（sensevoice）chunk 管理",
        ]),
        ("引擎層", 5.1, ACCENT2, [
            "Google STT — 雲端 chirp_3，asia-northeast1，gRPC 串流",
            "ElevenLabs Scribe v2 — 雲端 WebSocket RT，~150ms TTFT",
            "OpenAI Whisper — 本地端 large-v3 / turbo（faster-whisper）",
            "SenseVoice — 本地端離線，RTF≈0.1，情緒/事件辨識",
        ]),
    ]

    for label, y, color, items in layers:
        # 層標籤
        lbl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(1.8), Inches(1.5))
        lbl.fill.solid()
        lbl.fill.fore_color.rgb = color
        lbl.line.fill.background()
        txf = lbl.text_frame
        txf.paragraphs[0].text = label
        txf.paragraphs[0].font.size = Pt(16)
        txf.paragraphs[0].font.bold = True
        txf.paragraphs[0].font.color.rgb = WHITE
        txf.paragraphs[0].font.name = "Calibri"
        txf.paragraphs[0].alignment = PP_ALIGN.CENTER
        txf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 內容
        add_card(slide, 2.9, y, 9.5, 1.5)
        add_text_in_card(slide, 2.9, y, 9.5, 1.5, [
            (item, False, TEXT_DARK) for item in items
        ], font_size=12)

    # 箭頭連接
    for y in [2.8, 4.7]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.5), Inches(y), Inches(0.4), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEXT_MID
        arrow.line.fill.background()

    # ══════════════════════════════════════════════════════════════════════
    # Slide 4: 即時語音辨識模組
    # ══════════════════════════════════════════════════════════════════════
    slide = add_content_slide(prs, "即時語音辨識模組（app_api.py）")

    # 串流模式表格
    table_data = [
        ["模式", "引擎", "延遲", "存庫", "適用場景"],
        ["dual（推薦）", "Scribe RT + Google 批次", "~150ms 顯示\n15s 確認", "✅", "日常監聽 + 歸檔"],
        ["scribe_rt", "ElevenLabs Scribe v2 RT", "~150ms TTFT", "❌", "最低延遲即時字幕"],
        ["google_stream", "Google gRPC 串流", "~500ms", "✅ final", "高精度串流"],
        ["sensevoice_local", "SenseVoiceSmall（本地端）", "~3s chunk", "✅", "機密語音離線辨識"],
        ["batch", "Google / Scribe 批次", "15s", "✅", "向下相容"],
    ]

    rows, cols = len(table_data), len(table_data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.3), Inches(11.5), Inches(3.2))
    tbl = tbl_shape.table
    col_widths = [Inches(2.0), Inches(3.0), Inches(1.8), Inches(1.0), Inches(3.7)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    for r, row_data in enumerate(table_data):
        for c, cell_text in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = cell_text
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)
                para.font.name = "Calibri"
                if r == 0:
                    para.font.bold = True
                    para.font.color.rgb = WHITE
                else:
                    para.font.color.rgb = TEXT_DARK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)

    # WebSocket 端點說明
    add_card(slide, 0.8, 4.8, 11.5, 2.0)
    add_text_in_card(slide, 0.8, 4.8, 11.5, 2.0, [
        ("WebSocket 端點", True, PRIMARY),
        ("WS /ws/stream/{channel_id}?mode=dual&backend=google", False, TEXT_DARK),
        ("訊息類型: engine_info → partial（即時中間結果）→ transcript（確認）→ confirmed（Google 批次確認）", False, TEXT_MID),
        ("音訊格式: PCM 16-bit LE, 16kHz Mono | AudioBuffer: 15s（batch）/ 3s（sensevoice）", False, TEXT_MID),
        ("處理函式: _handle_dual_mode / _handle_scribe_rt_mode / _handle_google_stream_mode / _handle_sensevoice_local_mode / _handle_batch_mode", False, TEXT_MID),
    ], font_size=11)

    # ══════════════════════════════════════════════════════════════════════
    # Slide 5: 語音檔案辨識模組
    # ══════════════════════════════════════════════════════════════════════
    slide = add_content_slide(prs, "語音檔案辨識模組（app_dashboard.py 6 步驟流程）")

    steps = [
        ("Step 1", "選擇模型", "Google STT / Scribe / Whisper\nGemini / SenseVoice", PRIMARY),
        ("Step 2", "選擇子模型", "chirp_3 / scribe_v1\nlarge-v3 / SenseVoiceSmall", ACCENT),
        ("Step 3", "載入音檔", "上傳 wav/mp3/m4a/flac\n瀏覽伺服器目錄", ACCENT2),
        ("Step 4", "詞彙優化", "PhraseSet 提示 + 術語修正\n+ 同音字修正（170 條）", ACCENT3),
        ("Step 5", "音訊前處理", "Silero VAD 過濾\nDeepFilterNet 降噪", PRIMARY),
        ("Step 6", "彙整輸出", "依時間排序合併\n簡繁轉換 + 事件命名", ACCENT),
    ]

    for i, (step, title, desc, color) in enumerate(steps):
        col = i % 3
        row = i // 3
        x = 0.8 + col * 4.0
        y = 1.3 + row * 2.8

        # Step 標籤
        lbl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(3.5), Inches(0.5))
        lbl.fill.solid()
        lbl.fill.fore_color.rgb = color
        lbl.line.fill.background()
        txf = lbl.text_frame
        txf.paragraphs[0].text = f"{step}  {title}"
        txf.paragraphs[0].font.size = Pt(14)
        txf.paragraphs[0].font.bold = True
        txf.paragraphs[0].font.color.rgb = WHITE
        txf.paragraphs[0].font.name = "Calibri"
        txf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 描述卡片
        add_card(slide, x, y + 0.55, 3.5, 1.8)
        txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.7), Inches(3.1), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = desc
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = TEXT_DARK
        tf.paragraphs[0].font.name = "Calibri"

    # ══════════════════════════════════════════════════════════════════════
    # Slide 6: 辨識引擎比較
    # ══════════════════════════════════════════════════════════════════════
    slide = add_content_slide(prs, "辨識引擎比較（5 大引擎）")

    engines = [
        ["引擎", "類型", "模型", "延遲", "中文精度", "離線", "特色"],
        ["Google STT", "雲端", "chirp_3", "~500ms", "★★★★★", "❌", "PhraseSet、講者辨識、區域端點"],
        ["Scribe v2 RT", "雲端", "scribe_v2_realtime", "~150ms", "★★★★", "❌", "WebSocket、最低延遲、自動 VAD"],
        ["Whisper", "本地端", "large-v3 / turbo", "~4s/段", "★★★★", "✅", "faster-whisper 加速、initial_prompt"],
        ["SenseVoice", "本地端", "iic/SenseVoiceSmall", "~70ms / 10s", "★★★★", "✅", "情緒辨識、事件偵測（笑/哭/掌聲）"],
        ["Gemini", "雲端", "2.5-flash / 2.5-pro", "~2s/段", "★★★★★", "❌", "多模態、語意理解、自訂 prompt"],
    ]

    rows, cols = len(engines), len(engines[0])
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.5))
    tbl = tbl_shape.table
    col_widths = [Inches(1.5), Inches(1.0), Inches(1.8), Inches(1.5), Inches(1.3), Inches(0.8), Inches(4.4)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    for r, row_data in enumerate(engines):
        for c, cell_text in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = cell_text
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)
                para.font.name = "Calibri"
                if r == 0:
                    para.font.bold = True
                    para.font.color.rgb = WHITE
                else:
                    para.font.color.rgb = TEXT_DARK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY

    # 模組位置與工廠函式說明
    add_card(slide, 0.5, 5.2, 12.3, 1.8)
    add_text_in_card(slide, 0.5, 5.2, 12.3, 1.8, [
        ("STT 模組位置 — scripts/models/", True, ACCENT3),
        ("model_google_stt.py / model_scribe.py / model_whisper.py / model_sensevoice.py / model_gemini.py", False, TEXT_DARK),
        ("工廠函式: create_stt_model(backend='google'|'scribe'|'sensevoice') — app_api.py 統一建立 STT 實例", False, TEXT_MID),
    ], font_size=12)

    # ══════════════════════════════════════════════════════════════════════
    # Slide 7: 資料庫架構
    # ══════════════════════════════════════════════════════════════════════
    slide = add_content_slide(prs, "資料庫架構（SQLite WAL + FTS5）")

    # 即時 DB
    add_card(slide, 0.8, 1.3, 5.5, 3.0)
    add_text_in_card(slide, 0.8, 1.3, 5.5, 3.0, [
        ("即時串流表 — aiSpeechMulti.db", True, PRIMARY),
        ("", False, TEXT_MID),
        ("transcripts 表", True, ACCENT),
        ("id | channel_id | transcript | confidence", False, TEXT_DARK),
        ("stt_backend | use_vad | use_denoise | created_at", False, TEXT_DARK),
        ("", False, TEXT_MID),
        ("索引: idx_ch (channel_id), idx_ts (created_at)", False, TEXT_MID),
        ("WAL 模式 — 支援併發讀寫（多 WebSocket 連線）", False, TEXT_MID),
    ], font_size=12)

    # 批次 DB
    add_card(slide, 6.8, 1.3, 5.5, 3.0)
    add_text_in_card(slide, 6.8, 1.3, 5.5, 3.0, [
        ("批次辨識相關表（utils/db_manager.py）", True, PRIMARY),
        ("", False, TEXT_MID),
        ("events — 事件 (event_name, hazard_level 0-5, model_type)", True, ACCENT),
        ("audio_files — 音檔 (filename, archive_path, file_hash)", True, ACCENT2),
        ("transcriptions — 辨識結果 (transcript, status, FTS5 同步)", True, ACCENT3),
        ("keywords — 關鍵字 (keyword, hazard_level, source)", True, DANGER),
        ("", False, TEXT_MID),
        ("transcriptions_fts — FTS5 trigram tokenizer（中文子字串檢索）", False, TEXT_MID),
    ], font_size=12)

    # 資料流
    add_card(slide, 0.8, 4.6, 11.5, 2.5)
    add_text_in_card(slide, 0.8, 4.6, 11.5, 2.5, [
        ("資料流", True, PRIMARY),
        ("", False, TEXT_MID),
        ("即時串流: Browser PCM → WebSocket → STT → transcripts 表 → 前端推播", False, TEXT_DARK),
        ("批次辨識: 音檔上傳 → events → audio_files → STT → transcriptions → FTS5 同步", False, TEXT_DARK),
        ("關鍵字提取: transcriptions → Gemini → keywords → display.html 即時高亮", False, TEXT_DARK),
        ("事件歸檔: 音檔 → SHA-256 → data/audio_archive/ → audio_files.archive_path", False, TEXT_DARK),
    ], font_size=12)

    # ══════════════════════════════════════════════════════════════════════
    # Slide 8: 詞彙系統 + 音訊前處理
    # ══════════════════════════════════════════════════════════════════════
    slide = add_content_slide(prs, "三層詞彙優化系統 + 音訊前處理")

    # 三層詞彙
    vocab_layers = [
        ("第一層 — 辨識前提示", "Google PhraseSet\n(184 條術語)", "vocabulary/master_vocabulary.csv\nboost_value 權重排序，僅 Google STT 有效", PRIMARY),
        ("第二層 — 辨識後修正", "術語 + 同音字修正\n(~170 條規則)", "vocabulary/correction_dict.py\n月台編號、軍事數字讀法、設備代號修正", ACCENT),
        ("第三層 — 緊急關鍵字", "alert_keywords.json\n+ 危害等級分類", "vocabulary/alert_keywords.json\n0-5 級危害分級，display.html 即時高亮", ACCENT2),
    ]

    for i, (title, subtitle, desc, color) in enumerate(vocab_layers):
        x = 0.8 + i * 4.0
        lbl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.3), Inches(3.5), Inches(0.5))
        lbl.fill.solid()
        lbl.fill.fore_color.rgb = color
        lbl.line.fill.background()
        txf = lbl.text_frame
        txf.paragraphs[0].text = title
        txf.paragraphs[0].font.size = Pt(13)
        txf.paragraphs[0].font.bold = True
        txf.paragraphs[0].font.color.rgb = WHITE
        txf.paragraphs[0].font.name = "Calibri"
        txf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_card(slide, x, 1.85, 3.5, 2.0)
        add_text_in_card(slide, x, 1.85, 3.5, 2.0, [
            (subtitle, True, TEXT_DARK),
            ("", False, TEXT_MID),
            (desc, False, TEXT_MID),
        ], font_size=11)

    # 音訊前處理
    add_card(slide, 0.8, 4.2, 5.5, 2.8)
    add_text_in_card(slide, 0.8, 4.2, 5.5, 2.8, [
        ("音訊前處理管線（utils/）", True, PRIMARY),
        ("", False, TEXT_MID),
        ("1. DeepFilterNet 降噪 — noise_filter.py", False, TEXT_DARK),
        ("   AI 背景噪音抑制；無線電窄頻建議關閉", False, TEXT_MID),
        ("2. Silero VAD 靜音過濾 — vad_filter.py", False, TEXT_DARK),
        ("   門檻 0.0-1.0，全域單例 thread-safe", False, TEXT_MID),
        ("3. ffmpeg 格式轉換 → LINEAR16 PCM 16kHz Mono", False, TEXT_DARK),
        ("4. 動態設定 — POST /api/settings 即時生效", False, TEXT_DARK),
    ], font_size=12)

    # 詞彙表結構
    add_card(slide, 6.8, 4.2, 5.5, 2.8)
    add_text_in_card(slide, 6.8, 4.2, 5.5, 2.8, [
        ("master_vocabulary.csv 欄位（184 條）", True, PRIMARY),
        ("", False, TEXT_MID),
        ("term — 專業術語（OCC、月台門、CBTC、ATP...）", False, TEXT_DARK),
        ("category — equipment / location / action / personnel / numeric", False, TEXT_DARK),
        ("boost_value — 權重 0-20（影響辨識優先序）", False, TEXT_DARK),
        ("alert_level — 危害等級（0 正常 ~ 5 緊急）", False, TEXT_DARK),
        ("pinyin / common_error — 拼音與常見誤辨對照", False, TEXT_DARK),
    ], font_size=12)

    # ══════════════════════════════════════════════════════════════════════
    # Slide 9: 前端介面 + Streamlit 10 頁面
    # ══════════════════════════════════════════════════════════════════════
    slide = add_content_slide(prs, "前端介面 + Streamlit 管理面板（10 頁面）")

    # 前端 HTML
    add_card(slide, 0.8, 1.3, 5.5, 2.5)
    add_text_in_card(slide, 0.8, 1.3, 5.5, 2.5, [
        ("前端頁面（static/）", True, PRIMARY),
        ("", False, TEXT_MID),
        ("index.html — 單頻道串流控制 + 辨識顯示", False, TEXT_DARK),
        ("monitor.html — 6 路即時監控（自動輪詢）", False, TEXT_DARK),
        ("display.html — 全螢幕關鍵字投影（控制室）", False, TEXT_DARK),
        ("技術: Web Audio API + WebSocket + Vanilla JS", False, TEXT_MID),
    ], font_size=12)

    # Streamlit 10 頁面
    add_card(slide, 6.8, 1.3, 5.5, 2.5)
    add_text_in_card(slide, 6.8, 1.3, 5.5, 2.5, [
        ("Streamlit 10 頁面（app_dashboard.py）", True, PRIMARY),
        ("", False, TEXT_MID),
        ("home / monitor / speech / running / management", False, TEXT_DARK),
        ("search / stats / vocabulary / evaluation / offline_monitor", False, TEXT_DARK),
        ("功能: 即時監控、批次辨識、事件管理、全文檢索、", False, TEXT_MID),
        ("      統計圖表、詞彙線上編輯、CER/WER 評測、資料夾監看", False, TEXT_MID),
    ], font_size=12)

    # REST API
    add_card(slide, 0.8, 4.1, 11.5, 3.0)
    add_text_in_card(slide, 0.8, 4.1, 11.5, 3.0, [
        ("REST API 端點（FastAPI / app_api.py）", True, PRIMARY),
        ("", False, TEXT_MID),
        ("GET  /api/channels       — 6 路管道狀態（連線時間、辨識筆數、引擎、模式）", False, TEXT_DARK),
        ("GET  /api/transcripts    — 查詢辨識結果（分頁、頻道篩選、時間排序）", False, TEXT_DARK),
        ("GET  /api/settings       — 音訊前處理設定（VAD / 降噪 / 門檻）", False, TEXT_DARK),
        ("POST /api/settings       — 動態更新前處理設定（即時生效，無需重啟）", False, TEXT_DARK),
        ("GET  /api/keywords       — 關鍵字清單（display.html 即時比對高亮）", False, TEXT_DARK),
        ("GET  /api/test_scribe    — Scribe RT 連線診斷（合成音訊測試）", False, TEXT_DARK),
        ("GET  /api/health         — 健康檢查 + 版本 + DB 狀態", False, TEXT_DARK),
    ], font_size=11)

    # ══════════════════════════════════════════════════════════════════════
    # Slide 10: 離線模式規劃（NEW）
    # ══════════════════════════════════════════════════════════════════════
    slide = add_content_slide(prs, "離線模式規劃（Mac mini M4 24GB）")

    # 系統架構
    add_card(slide, 0.8, 1.3, 5.5, 3.0)
    add_text_in_card(slide, 0.8, 1.3, 5.5, 3.0, [
        ("系統架構", True, PRIMARY),
        ("", False, TEXT_MID),
        ("輸入: IMA-ADPCM 無線電音檔", False, TEXT_DARK),
        ("→ ffmpeg 轉換 → 16kHz Mono WAV", False, TEXT_MID),
        ("", False, TEXT_MID),
        ("雙引擎並行（CPU 推論）:", True, ACCENT),
        ("• faster-whisper（CTranslate2 加速）", False, TEXT_DARK),
        ("• SenseVoiceSmall（FunASR）", False, TEXT_DARK),
        ("", False, TEXT_MID),
        ("輸出: Streamlit Web UI（port 8501）", False, TEXT_DARK),
        ("情緒標籤 + CSV / SRT 匯出", False, TEXT_MID),
    ], font_size=12)

    # 4 大功能模式
    add_card(slide, 6.8, 1.3, 5.5, 3.0)
    add_text_in_card(slide, 6.8, 1.3, 5.5, 3.0, [
        ("4 大功能模式", True, PRIMARY),
        ("", False, TEXT_MID),
        ("1. 單檔辨識 — 上傳音檔即時處理", False, TEXT_DARK),
        ("2. 資料夾監看 — 近即時自動處理新檔", False, TEXT_DARK),
        ("3. 批次處理 — 整批音檔一次辨識", False, TEXT_DARK),
        ("4. 雙引擎比對 — 同檔案兩引擎並陳", False, TEXT_DARK),
        ("", False, TEXT_MID),
        ("詞彙注入: faster-whisper prompt", False, TEXT_MID),
        ("           SenseVoice hotwords", False, TEXT_MID),
    ], font_size=12)

    # 效能基準與安全
    add_card(slide, 0.8, 4.6, 11.5, 2.5)
    add_text_in_card(slide, 0.8, 4.6, 11.5, 2.5, [
        ("效能基準（Mac mini M4 CPU）與安全設計", True, PRIMARY),
        ("", False, TEXT_MID),
        ("SenseVoiceSmall — 10s 音檔約 1s 推論（RTF 0.1x）  ★★★★★ 中文精度", False, TEXT_DARK),
        ("faster-whisper large-v3 — 10s 音檔約 15s 推論（RTF 1.5x）  ★★★★ 多語言", False, TEXT_DARK),
        ("", False, TEXT_MID),
        ("安全: 100% 離線（無雲端傳輸）| 網路隔離 | 防火牆防護 | USB 模型更新", False, TEXT_DARK),
        ("文件: 離線模式規劃/DEPLOYMENT_GUIDE.md", False, TEXT_MID),
    ], font_size=12)

    # ══════════════════════════════════════════════════════════════════════
    # Slide 11: 結尾
    # ══════════════════════════════════════════════════════════════════════
    slide = add_dark_slide(prs,
        "aiSpeechMulti v2.1",
        "多引擎即時語音辨識平台"
    )
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(11), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    items = [
        "GitHub: github.com/datadigshawn/aiSpeechMulti",
        "啟動 API: uvicorn app_api:app --host 0.0.0.0 --port 8000",
        "啟動面板: streamlit run app_dashboard.py",
        "OpenAPI 文件: http://localhost:8000/docs",
    ]
    for i, text in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xA0, 0xAE, 0xC0)
        p.font.name = "Calibri"
        p.space_before = Pt(8)

    # ── 儲存 ──────────────────────────────────────────────────────────────
    output_path = Path(__file__).parent / "aiSpeechMulti_Architecture.pptx"
    prs.save(str(output_path))
    print(f"✅ 簡報已產生: {output_path}")
    return output_path


if __name__ == "__main__":
    build_presentation()
